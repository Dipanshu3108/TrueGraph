"""Tests for PPTXParser."""

import os
import tempfile
import unittest

from libs.UDP.universal_document_parser.parsers.pptx_parser import PPTXParser


try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


@unittest.skipUnless(HAS_PPTX, "python-pptx not installed")
class TestPPTXParser(unittest.TestCase):
    def test_parse_pptx(self):
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            textbox = slide.shapes.add_textbox(0, 0, 100, 100)
            textbox.text_frame.text = "Slide content"
            prs.save(tmp.name)
            tmp_path = tmp.name

        try:
            parser = PPTXParser()
            result = parser.parse(tmp_path)
            self.assertEqual(len(result), 1)
            self.assertIn("Slide content", result[0]["content"])
        finally:
            os.unlink(tmp_path)


if __name__ == "__main__":
    unittest.main()
