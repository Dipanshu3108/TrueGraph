"""PPTX (and converted PPT) parser."""

import os
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional

from ..cleanup.text_cleanup import cleanup_page
from ..config import DEFAULT_CONFIG
from ..exceptions import DependencyError, ParseError
from ..logger import get_logger

logger = get_logger()

try:
    import pptx
    from pptx import Presentation
except ImportError as exc:  # pragma: no cover
    pptx = None  # type: ignore
    Presentation = None  # type: ignore


def _extract_shape_text(shape: Any) -> List[str]:
    """Recursively extract text from a shape and its grouped children."""
    texts: List[str] = []
    if hasattr(shape, "text") and shape.text:
        texts.append(shape.text.strip())
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            texts.extend(_extract_shape_text(child))
    if hasattr(shape, "table"):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                texts.append(" | ".join(cells))
    return [t for t in texts if t]


class PPTXParser:
    """Parse PPTX files and converted PPT files."""

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        self.config = {**DEFAULT_CONFIG, **(config or {})}

    def parse(self, path: str) -> List[Dict[str, Any]]:
        if Presentation is None:
            raise DependencyError(
                "python-pptx is required. Install it with: pip install python-pptx"
            )

        try:
            prs = Presentation(path)
        except Exception as exc:
            raise ParseError(f"Failed to open PPTX {path}: {exc}") from exc

        pages: List[Dict[str, Any]] = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []

            # Shape text (including grouped shapes and tables)
            for shape in slide.shapes:
                texts.extend(_extract_shape_text(shape))

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    texts.append(f"Notes: {notes}")

            content = cleanup_page("\n".join(texts))
            pages.append({"page_no": idx, "content": content})

        return pages

    def convert_ppt_to_pptx(self, path: str) -> str:
        """Convert a .ppt file to .pptx using LibreOffice headless.

        Returns the path to the generated .pptx file.
        """
        if not self.config.get("enable_ppt_conversion", True):
            raise DependencyError(
                "PPT conversion is disabled. Enable it or convert the file manually."
            )

        libreoffice = self._find_libreoffice()
        if libreoffice is None:
            raise DependencyError(
                "LibreOffice is required to convert .ppt files. "
                "Install LibreOffice and ensure 'soffice' or 'libreoffice' is on PATH."
            )

        path_obj = Path(path).resolve()
        out_dir = tempfile.mkdtemp(prefix="udp_ppt_convert_")
        try:
            subprocess.run(
                [
                    libreoffice,
                    "--headless",
                    "--convert-to",
                    "pptx",
                    "--outdir",
                    out_dir,
                    str(path_obj),
                ],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
            )
        except subprocess.CalledProcessError as exc:
            raise ParseError(f"LibreOffice failed to convert {path}: {exc.stderr}") from exc

        expected_name = path_obj.stem + ".pptx"
        out_path = Path(out_dir) / expected_name
        if not out_path.exists():
            raise ParseError(f"LibreOffice did not produce expected output: {out_path}")
        return str(out_path)

    @staticmethod
    def _find_libreoffice() -> Optional[str]:
        """Locate the LibreOffice/soffice executable."""
        candidates = ["soffice", "libreoffice"]
        if os.name == "nt":
            candidates = ["soffice.exe", "libreoffice.exe", r"C:\Program Files\LibreOffice\program\soffice.exe"]

        for candidate in candidates:
            if os.path.isabs(candidate) and os.path.isfile(candidate):
                return candidate
            found = _which(candidate)
            if found:
                return found
        return None


def _which(cmd: str) -> Optional[str]:
    """Portable 'which' implementation."""
    for path in os.environ.get("PATH", "").split(os.pathsep):
        full = os.path.join(path, cmd)
        if os.path.isfile(full):
            return full
    return None
